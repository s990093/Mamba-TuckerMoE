#!/usr/bin/env python3
"""
Export Presentation/index.html → high-resolution PPTX
Uses Chrome headless to screenshot each Reveal.js slide, then builds PPTX.

Usage:
  python3 export_pptx.py                    # outputs presentation.pptx
  python3 export_pptx.py --slides 13        # override slide count
  python3 export_pptx.py --scale 2          # retina 2× (default)
"""

import argparse, http.server, json, os, pathlib, shutil, subprocess, sys
import tempfile, threading, time

# ─── constants ──────────────────────────────────────────────────────────────
CHROME = "/Applications/Google Chrome.app/Contents/MacOS/Google Chrome"
SLIDE_W_PT  = 1280   # logical width  (pixels at scale=1)
SLIDE_H_PT  = 720    # logical height
PORT        = 59127
HERE        = pathlib.Path(__file__).parent
HTML_FILE   = HERE / "index.html"
OUT_PPTX    = HERE / "presentation.pptx"
SLIDE_COUNT = 13     # default

# ─── argument parsing ────────────────────────────────────────────────────────
parser = argparse.ArgumentParser()
parser.add_argument("--slides", type=int, default=SLIDE_COUNT)
parser.add_argument("--scale",  type=int, default=2,
                    help="Device pixel ratio (1=1280×720, 2=2560×1440 retina)")
args = parser.parse_args()

SCALE = args.scale
W_PX  = SLIDE_W_PT * SCALE
H_PX  = SLIDE_H_PT * SCALE

# ─── serve the HTML locally ──────────────────────────────────────────────────
class _SilentHandler(http.server.SimpleHTTPRequestHandler):
    def log_message(self, *_): pass

server = http.server.HTTPServer(("127.0.0.1", PORT), _SilentHandler)
server.allow_reuse_address = True
server_thread = threading.Thread(target=server.serve_forever, daemon=True)
os.chdir(HERE)
server_thread.start()
print(f"[server] http://127.0.0.1:{PORT}/index.html")
time.sleep(0.5)  # let OS bind

# ─── capture slides via Chrome headless + DevTools Protocol ─────────────────
tmp = pathlib.Path(tempfile.mkdtemp(prefix="pptx_slides_"))
print(f"[export] capturing {args.slides} slides at {W_PX}×{H_PX} px …")

# Build a single JS string that visits each slide and screenshots it,
# passed via --run-all-compositor-stages-before-draw + remote debugging.
# We drive it with Chrome's --remote-debugging-port and a helper script.

helper_js = tmp / "capture.js"
slides_dir = tmp / "slides"
slides_dir.mkdir()

helper_js.write_text(f"""
const http = require('http');
const fs   = require('fs');
const path = require('path');

const BASE_URL = 'http://127.0.0.1:{PORT}/index.html';
const SLIDES   = {args.slides};
const OUT_DIR  = {json.dumps(str(slides_dir))};
const DELAY    = 900;  // ms after slide transition

// --- minimal CDP client ---
function cdpRequest(port, method, params, cb) {{
  const body = JSON.stringify({{ id: 1, method, params: params || {{}} }});
  const req = http.request({{ host:'127.0.0.1', port, path:'/json/new', method:'GET' }}, () => {{}});
  req.on('error', () => {{}});
  req.end();
  // we'll use the websocket approach via a separate ws call instead
  cb && cb();
}}

// Helper: sleep
const sleep = ms => new Promise(r => setTimeout(r, ms));

async function run() {{
  // Connect to Chrome devtools
  const listResp = await fetch(`http://127.0.0.1:9222/json`);
  const tabs = await listResp.json();
  const tab = tabs.find(t => t.url && t.url.includes('index.html')) || tabs[0];
  if (!tab) {{ console.error('No tab found'); process.exit(1); }}

  const ws = require('ws');
  const client = new ws(tab.webSocketDebuggerUrl);

  let msgId = 1;
  const pending = {{}};
  client.on('message', raw => {{
    const msg = JSON.parse(raw);
    if (msg.id && pending[msg.id]) pending[msg.id](msg.result);
  }});

  const send = (method, params) => new Promise(resolve => {{
    const id = msgId++;
    pending[id] = resolve;
    client.send(JSON.stringify({{ id, method, params: params || {{}} }}));
  }});

  await new Promise(r => client.on('open', r));

  // Enable domains
  await send('Runtime.enable');
  await send('Page.enable');

  // Navigate to presentation
  await send('Page.navigate', {{ url: BASE_URL + '#/0' }});
  await sleep(2000);  // wait for Reveal.js to init + MathJax

  // Set viewport
  await send('Emulation.setDeviceMetricsOverride', {{
    width: {SLIDE_W_PT}, height: {SLIDE_H_PT},
    deviceScaleFactor: {SCALE}, mobile: false
  }});

  for (let i = 0; i < SLIDES; i++) {{
    // Navigate to slide i
    await send('Runtime.evaluate', {{
      expression: `Reveal.slide(${{i}}, 0)`,
      awaitPromise: false
    }});
    await sleep(DELAY);

    // Hide nav controls + progress for clean screenshot
    await send('Runtime.evaluate', {{
      expression: `
        document.querySelector('.reveal .controls') && (document.querySelector('.reveal .controls').style.display='none');
        document.querySelector('.reveal .progress') && (document.querySelector('.reveal .progress').style.display='none');
      `
    }});
    await sleep(100);

    const shot = await send('Page.captureScreenshot', {{
      format: 'png',
      clip: {{ x:0, y:0, width:{SLIDE_W_PT}, height:{SLIDE_H_PT}, scale:{SCALE} }}
    }});
    const buf = Buffer.from(shot.data, 'base64');
    const out = path.join(OUT_DIR, `slide_${{String(i).padStart(3,'0')}}.png`);
    fs.writeFileSync(out, buf);
    console.log(`  captured slide ${{i+1}}/${{SLIDES}} → ${{out}}`);
  }}

  client.close();
  console.log('DONE');
  process.exit(0);
}}

run().catch(e => {{ console.error(e); process.exit(1); }});
""", encoding="utf-8")

# Check if Node.js ws module is available; if not use a pure-Python CDP approach
def check_node_ws():
    r = subprocess.run(["node", "-e", "require('ws')"], capture_output=True)
    return r.returncode == 0

use_node = check_node_ws()

# ─── Python-native CDP approach (no node ws needed) ──────────────────────────
import socket, struct, base64, hashlib, select

def ws_handshake(sock, host, path):
    key = base64.b64encode(os.urandom(16)).decode()
    req = (f"GET {path} HTTP/1.1\r\n"
           f"Host: {host}\r\n"
           f"Upgrade: websocket\r\n"
           f"Connection: Upgrade\r\n"
           f"Sec-WebSocket-Key: {key}\r\n"
           f"Sec-WebSocket-Version: 13\r\n\r\n")
    sock.sendall(req.encode())
    resp = b""
    while b"\r\n\r\n" not in resp:
        resp += sock.recv(4096)

def ws_send(sock, data: str):
    payload = data.encode()
    n = len(payload)
    mask = os.urandom(4)
    if n < 126:
        header = bytes([0x81, 0x80 | n]) + mask
    elif n < 65536:
        header = bytes([0x81, 0xFE]) + struct.pack(">H", n) + mask
    else:
        header = bytes([0x81, 0xFF]) + struct.pack(">Q", n) + mask
    masked = bytes(b ^ mask[i % 4] for i, b in enumerate(payload))
    sock.sendall(header + masked)

def ws_recv(sock):
    def _recv_exact(n):
        buf = b""
        while len(buf) < n:
            chunk = sock.recv(n - len(buf))
            if not chunk:
                raise ConnectionError("socket closed")
            buf += chunk
        return buf
    b0, b1 = _recv_exact(2)
    masked = (b1 & 0x80) != 0
    plen = b1 & 0x7F
    if plen == 126:
        plen = struct.unpack(">H", _recv_exact(2))[0]
    elif plen == 127:
        plen = struct.unpack(">Q", _recv_exact(8))[0]
    mask_key = _recv_exact(4) if masked else b""
    data = _recv_exact(plen)
    if masked:
        data = bytes(b ^ mask_key[i % 4] for i, b in enumerate(data))
    return json.loads(data.decode())

def launch_chrome_debug():
    user_data = tmp / "chrome_profile"
    user_data.mkdir(exist_ok=True)
    # Kill any previous instance on port 9222
    subprocess.run(["pkill", "-f", "remote-debugging-port=9222"],
                   capture_output=True)
    time.sleep(1)
    cmd = [
        CHROME,
        "--headless=new",
        "--disable-gpu",
        "--no-sandbox",
        "--disable-dev-shm-usage",
        "--disable-extensions",
        "--no-first-run",
        "--no-default-browser-check",
        "--disable-background-networking",
        f"--remote-debugging-port=9222",
        f"--user-data-dir={user_data}",
        f"--window-size={SLIDE_W_PT},{SLIDE_H_PT}",
        "about:blank",
    ]
    proc = subprocess.Popen(cmd, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
    return proc

def get_ws_url(port=9222, retries=20):
    import urllib.request
    for i in range(retries):
        try:
            with urllib.request.urlopen(
                    f"http://127.0.0.1:{port}/json", timeout=2) as r:
                tabs = json.loads(r.read())
                for t in tabs:
                    if t.get("type") == "page":
                        return t["webSocketDebuggerUrl"]
        except Exception:
            pass
        time.sleep(1)
        if i == 4:
            print("    [waiting for Chrome to start…]")
    raise RuntimeError("Could not connect to Chrome DevTools after 20 s")

def capture_slides_python(n_slides, scale, out_dir):
    chrome_proc = launch_chrome_debug()
    try:
        ws_url = get_ws_url()
        # parse ws url: ws://127.0.0.1:9222/devtools/page/...
        import urllib.parse
        parsed = urllib.parse.urlparse(ws_url)
        sock = socket.create_connection((parsed.hostname, parsed.port), timeout=10)
        ws_handshake(sock, f"{parsed.hostname}:{parsed.port}", parsed.path)

        msg_id = [1]
        def send(method, params=None):
            mid = msg_id[0]; msg_id[0] += 1
            ws_send(sock, json.dumps({"id": mid, "method": method, "params": params or {}}))
            # read until we get the response for this id
            while True:
                msg = ws_recv(sock)
                if msg.get("id") == mid:
                    return msg.get("result", {})

        # Enable
        send("Runtime.enable")
        send("Page.enable")

        # Navigate to presentation
        target_url = f"http://127.0.0.1:{PORT}/index.html"
        print(f"    navigating to {target_url} …")
        send("Page.navigate", {"url": target_url})
        # Wait for Reveal.js + MathJax to fully initialize
        time.sleep(5)

        # Set exact viewport
        send("Emulation.setDeviceMetricsOverride", {
            "width": SLIDE_W_PT, "height": SLIDE_H_PT,
            "deviceScaleFactor": scale, "mobile": False
        })
        time.sleep(0.5)

        for i in range(n_slides):
            # Go to slide
            send("Runtime.evaluate", {
                "expression": f"Reveal.slide({i}, 0)",
                "awaitPromise": False
            })
            time.sleep(1.0)

            # Hide controls & progress bar
            send("Runtime.evaluate", {
                "expression": (
                    "['controls','progress'].forEach(c=>{"
                    "  var el=document.querySelector('.reveal .'+c);"
                    "  if(el) el.style.display='none';"
                    "});"
                )
            })
            time.sleep(0.15)

            result = send("Page.captureScreenshot", {
                "format": "png",
                "clip": {
                    "x": 0, "y": 0,
                    "width": SLIDE_W_PT, "height": SLIDE_H_PT,
                    "scale": scale
                }
            })
            img_data = base64.b64decode(result["data"])
            out_path = out_dir / f"slide_{i:03d}.png"
            out_path.write_bytes(img_data)
            print(f"  slide {i+1:2d}/{n_slides}  →  {out_path.name}  ({len(img_data)//1024} KB)")

        sock.close()
    finally:
        chrome_proc.terminate()
        time.sleep(1)

# ─── capture ─────────────────────────────────────────────────────────────────
capture_slides_python(args.slides, SCALE, slides_dir)
server.shutdown()

# ─── assemble PPTX ───────────────────────────────────────────────────────────
print(f"\n[pptx] assembling {args.slides} slides …")
try:
    from pptx import Presentation as Prs
    from pptx.util import Pt, Emu
    from pptx.enum.text import PP_ALIGN
    import io

    prs = Prs()
    prs.slide_width  = Emu(int(SLIDE_W_PT * 914400 / 96))  # 1280px @ 96dpi
    prs.slide_height = Emu(int(SLIDE_H_PT * 914400 / 96))  # 720px @ 96dpi

    blank_layout = prs.slide_layouts[6]  # completely blank

    pngs = sorted(slides_dir.glob("slide_*.png"))
    for png in pngs:
        slide = prs.slides.add_slide(blank_layout)
        pic = slide.shapes.add_picture(
            str(png),
            left=0, top=0,
            width=prs.slide_width,
            height=prs.slide_height
        )

    prs.save(str(OUT_PPTX))
    print(f"[pptx] saved → {OUT_PPTX}  ({OUT_PPTX.stat().st_size//1024} KB)")

except ImportError as e:
    print(f"[pptx] python-pptx not available ({e}); slides saved as PNG in {slides_dir}")

# ─── also export a high-res PDF ──────────────────────────────────────────────
try:
    from PIL import Image as PILImage

    out_pdf = HERE / "presentation.pdf"
    imgs = []
    for png in sorted(slides_dir.glob("slide_*.png")):
        img = PILImage.open(png).convert("RGB")
        imgs.append(img)

    if imgs:
        imgs[0].save(
            str(out_pdf), "PDF", resolution=144,
            save_all=True, append_images=imgs[1:]
        )
        print(f"[pdf]  saved → {out_pdf}  ({out_pdf.stat().st_size//1024} KB)")
except ImportError:
    print("[pdf]  Pillow not available; skipping PDF export")

# ─── cleanup temp ────────────────────────────────────────────────────────────
shutil.rmtree(tmp, ignore_errors=True)
print("\nDone.")
